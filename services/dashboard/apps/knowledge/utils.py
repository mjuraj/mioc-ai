import mimetypes
import os
import re
from io import BytesIO

import docx2txt
import openpyxl
import pptx
from PyPDF2 import PdfReader
from django.conf import settings
from forge_storage import AzureStorageHandler


def strip_extension(filename):
    return os.path.splitext(filename)[0]


def upload_txt_file_to_azure(file_name, text):
    azure_storage = AzureStorageHandler(
        settings.AZURE_ACCOUNT_URL, settings.AZURE_ACCOUNT_KEY, settings.AZURE_MEDIA_CONTAINER)
    azure_storage.upload(file_name, text)


def remove_file_from_azure(file_name):
    azure_storage = AzureStorageHandler(
        settings.AZURE_ACCOUNT_URL, settings.AZURE_ACCOUNT_KEY, settings.AZURE_MEDIA_CONTAINER)
    azure_storage.delete(file_name)


def extract_text_from_file(file: bytes, filepath: str) -> str:
    mimetype, _ = mimetypes.guess_type(filepath)
    if mimetype == "application/pdf":
        # Extract text from pdf using PyPDF2
        reader = PdfReader(BytesIO(file))
        extracted_text = " ".join([page.extract_text().replace('\n', ' ') for page in reader.pages])
    elif mimetype == "text/plain" or mimetype == "text/markdown" or mimetype == "text/csv":
        # Read text from plain text file
        try:
            data = file.decode('utf-8')
        except:
            data = file.decode('utf-8', errors='ignore')
            # Remove non UTF-8 characters
            data = re.sub('[^\x00-\x7F]+', '', data)
        extracted_text = data
    elif mimetype == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(BytesIO(file))
    elif mimetype == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    elif mimetype == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        extracted_text = ""
        workbook = openpyxl.load_workbook(filename=BytesIO(file))
        sheets = workbook.sheetnames
        for sheet_name in sheets:
            worksheet = workbook[sheet_name]
            for row in worksheet.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        extracted_text += str(cell) + " "
                extracted_text += "\n"
    else:
        # Unsupported file type
        raise ValueError("Unsupported file type: {}".format(mimetype))

    return extracted_text


class ChunkedFileUploadHandler:
    temp_directory = os.path.join(settings.MEDIA_ROOT, 'temp')

    def __init__(self, filename):
        self.filename = filename
        self.chunk_path = os.path.join(self.temp_directory, f'{self.filename}')

    def save_chunk(self, chunk):
        if not os.path.exists(self.temp_directory):
            os.makedirs(self.temp_directory)
        with open(self.chunk_path, 'ab+') as destination:
            destination.write(chunk)

    def get_whole_file(self):
        with open(self.chunk_path, 'rb') as destination:
            return destination.read()

    def delete_temp_file(self):
        os.remove(self.chunk_path)
